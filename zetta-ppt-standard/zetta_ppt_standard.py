# -*- coding: utf-8 -*-
"""
zetta_ppt_standard.py  —  롯데마트 온라인/ZETTA 표준 PPT 빌더 (v4)
================================================================
'_양식__내부토의_공통양식_v3_' 파일과 '11차 Ocado 정기협의체' 산출물을
역설계하여 확정한 디자인 토큰·슬라이드 원형을 코드로 encode.

설계 원칙
  - 모든 좌표는 cm (실측 기준값). EMU 변환은 Cm() 이 처리.
  - 폰트는 맑은 고딕 3중 지정(latin·ea·cs) — v4 확정. (v3 테마 latin=Arial 을 교정)
  - 슬라이드 크롬(헤더/단위/페이지번호/각주)은 좌표 고정 → 전 페이지 정합.
  - 원형(archetype) 함수는 slide 를 반환 → 호출측에서 본문 자유 배치.
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy
import re

# ─────────────────────────────────────────────────────────────
# 1. 디자인 토큰 (역설계 확정값)
# ─────────────────────────────────────────────────────────────
PAGE_W_CM, PAGE_H_CM = 27.52, 19.05          # A4 가로 (9906000 x 6858000 EMU)
FONT       = "맑은 고딕"                       # 단일 폰트 — 한글·영문·숫자 공통 3중 지정 (v4.1)

# 팔레트 (theme1.xml 실측)
NAVY       = RGBColor(0x00, 0x00, 0x66)       # accent1  면 채움 전용(결론 박스 등) — 글자색 금지(v4.1)
CRIMSON    = RGBColor(0xC3, 0x0C, 0x3E)       # accent3  경고·리스크 강조
BLACK      = RGBColor(0x22, 0x22, 0x22)       # [P계열] 제목
INK        = RGBColor(0x00, 0x00, 0x00)       # 본문 기본
RED_GUIDE  = RGBColor(0xFF, 0x00, 0x00)       # [확인] 가이드(최종본 삭제 대상)
TH_PRIMARY = RGBColor(0xEA, 0xEE, 0xF6)       # 표 헤더 1차
TH_SECOND  = RGBColor(0xD8, 0xE0, 0xEC)       # 표 헤더 2차
TH_DARK    = RGBColor(0x00, 0x1E, 0x62)       # 강조 헤더(짙은 네이비)
GRID       = RGBColor(0xBF, 0xBF, 0xBF)       # 표 격자
SUB_GRAY   = RGBColor(0x59, 0x59, 0x59)       # 보조 텍스트·시사점 탭 회색
HL_FILL    = RGBColor(0xDE, 0xEB, 0xF7)       # 표 강조 행 연블루 (v4.1 산출물 실측)
HL_TEXT    = RGBColor(0x00, 0x20, 0x60)       # 표 강조 행 텍스트 — 파란 글자 금지의 유일 예외
DIVIDER_BG = RGBColor(0xF2, 0xF2, 0xF2)       # 챕터 간지 회색 풀블리드(bg1 lumMod95%)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

# 크롬 좌표 (실측 기반, v4.1 여백 +0.5mm 조정)  L, T, W, H  (cm)
# 여백·충전 규칙(v4.2): 좌·우·하단 1.5cm 고정 '액자' 여백. 그 안쪽(BODY_W × CONTENT_BOTTOM)은
# 빈 곳 없이 충전 — 객체를 세로로 분배해 마지막 객체(각주 포함) 하단선이 CONTENT_BOTTOM 에 닿게.
MARGIN_L, MARGIN_R = 1.50, 1.50               # 좌·우 여백 1.5cm
MARGIN_B = 1.50                               # 하단 여백 1.5cm
BODY_W = PAGE_W_CM - MARGIN_L - MARGIN_R       # 본문 폭 24.52
CONTENT_BOTTOM = PAGE_H_CM - MARGIN_B          # 콘텐츠 하한선 17.55 (여기까지 채움)
HDR   = (MARGIN_L, 1.20, BODY_W, 0.90)         # 좌상단 브래킷 헤더 (좌여백 1.5)
LEAD  = (MARGIN_L, 2.90, 21.50, 0.65)          # 제목 하단 리드메시지 (■ 한 줄)
UNIT  = (PAGE_W_CM - MARGIN_R - 3.30, 2.35, 3.30, 0.60)   # 우상단 단위 (우여백 1.5 정렬)
FOOT  = (MARGIN_L, CONTENT_BOTTOM - 0.55, BODY_W, 0.55)   # 좌하단 각주 (하단선 = CONTENT_BOTTOM)
BODY_TOP = 3.00                               # 본문 시작 y (리드메시지 미사용 시)

# 수직 간격 토큰 (정기협의체 산출물 실측 — 개체 간 간격 표준)
BODY_TOP_LEAD = 4.10                          # 본문 시작 y (리드메시지 사용 시, 리드 하단 +0.55)
GAP_SECTION   = 0.50                          # 본문 블록(소제목 묶음) 종료 → 다음 소제목 간격
GAP_HEAD      = 0.55                          # 소제목 시작 → 첫 본문 요소 시작 (피치)
LINE_H        = 0.48                          # 개조식 불릿 행간 (피치)

# 개조식 불릿 체계 (정기협의체 실측): ■ 리드(크롬) → • 주항목(볼드) → - 하위 → · 세부
BULLET_MARKS  = ("•", "-", "·")
BULLET_INDENT = 0.55                          # 하위 수준당 들여쓰기 (cm)

# 본장 2단 컬럼 프레임 (정기협의체 8·14/15 실측): 헤더 바 + 중앙 세로선 + 결론 박스
COL_GAP      = 0.66                           # 컬럼 간 간격
COL_HEADER_H = 0.75                           # 컬럼 헤더 바 높이
CONCL_H      = 0.90                           # 컬럼 하단 결론 박스 높이

# 제목 크기 (tier 별) — v4.1 교정: v3 실측(20pt 단일)로 회귀, tier 는 색상만 구분
TITLE_PT = {"P": 20, "F": 20}                 # P계열 검정 / F계열 네이비, 공통 20pt

# 표준 타이포 스케일 (v4.1, 정기협의체 실측) — 임의 크기 사용 금지
FONT_PT = {
    "title":    20,   # 브래킷 헤더 (볼드 검정 #222222, P/F 공통)
    "lead":     14,   # ■ 리드메시지 (볼드 검정)
    "bullet0":  12,   # • 주항목 (볼드 검정)
    "bullet1":  11,   # - 하위 (일반 검정)
    "bullet2":  10,   # · 세부 (일반 검정)
    "table":    10,   # 표 본문 (헤더: 검정 볼드 동일 크기)
    "unit":     10,   # 우상단 단위 (볼드)
    "footnote":  9,   # 좌하단 각주
}


# ─────────────────────────────────────────────────────────────
# 2. 저수준 헬퍼
# ─────────────────────────────────────────────────────────────
PAREN_DROP_PT = 1.0                              # 본문 이어지는 괄호는 한 사이즈 작게
_PAREN_RE = re.compile(r"\([^()]*\)")


def _emit_runs(p, text, size, bold, color, name=FONT):
    """문단에 run 생성 — 본문 이어지는 `(…)`는 한 사이즈 작게 + 여는 괄호 앞 한 칸 강제 (v4.1)."""
    text = re.sub(r"(?<=\S)\(", " (", text)      # 본문↔괄호 사이 한 칸
    text = re.sub(r" {2,}\(", " (", text)        # 이중 공백 방지
    small = None if size is None else max(size - PAREN_DROP_PT, 1)
    idx = 0
    for m in _PAREN_RE.finditer(text):
        if m.start() > idx:
            r = p.add_run(); r.text = text[idx:m.start()]
            _set_font(r, size, bold, color, name)
        r = p.add_run(); r.text = m.group()
        _set_font(r, small, bold, color, name)
        idx = m.end()
    if idx < len(text) or not p.runs:
        r = p.add_run(); r.text = text[idx:]
        _set_font(r, size, bold, color, name)


def _set_font(run, size=None, bold=None, color=None, name=FONT):
    """맑은 고딕 3중 지정(latin·ea·cs) + 속성. v4 폰트 규칙."""
    f = run.font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    f.name = name                              # latin
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):               # ea + cs 를 명시적으로 동일 지정
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def _txt(slide, l, t, w, h, text="", size=11, bold=False, color=INK,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name=FONT, wrap=True):
    """텍스트 박스 1개 생성 → shape 반환. 내부 패딩 0."""
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        _emit_runs(p, line, size, bold, color, name)
    return tb


def _cell_border(cell, color=None, w_pt=0.75):
    """셀 4방(lnL·lnR·lnT·lnB) 실선 격자 — 스키마 순서 보존 위해 tcPr 선두 삽입."""
    hexval = str(color) if color is not None else str(GRID)
    tcPr = cell._tc.get_or_add_tcPr()
    for tag in ("a:lnB", "a:lnT", "a:lnR", "a:lnL"):     # insert(0) 역순 → 최종 L,R,T,B
        el = tcPr.find(qn(tag))
        if el is not None:
            tcPr.remove(el)
        ln = tcPr.makeelement(qn(tag), {"w": str(int(w_pt * 12700)), "cap": "flat"})
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": hexval})
        fill.append(clr)
        ln.append(fill)
        tcPr.insert(0, ln)


def _line(slide, x1, y1, x2, y2, color=None, w_pt=0.75):
    """직선 커넥터 — 그림자 없음·실선 강제 (v4.1: 선에 그림자·점선 금지)."""
    ln = slide.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    ln.shadow.inherit = False                        # 그림자 제거
    lf = ln.line
    lf.color.rgb = color if color is not None else INK
    lf.width = Pt(w_pt)
    d = lf._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "solid"})
    lf._get_or_add_ln().append(d)                    # 실선 강제 (점선 방지)
    return ln


def _rect(slide, l, t, w, h, fill=None, line=None, line_w=0.75):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Cm(l), Cm(t), Cm(w), Cm(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


# ─────────────────────────────────────────────────────────────
# 3. 슬라이드 크롬 (전 페이지 공통)
# ─────────────────────────────────────────────────────────────
def add_header(slide, category, title, tier="P"):
    """좌상단 헤더. category 지정 시 `[분류] 제목`, None 이면 `■ 제목`
    (v4.1 — 원페이지 등 단독 보고서는 브래킷 분류 생략)."""
    l, t, w, h = HDR
    color = BLACK
    size = TITLE_PT[tier]
    tb = slide.shapes.add_textbox(Cm(l), Cm(t), Cm(w), Cm(h))
    tf = tb.text_frame; tf.word_wrap = True
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, 0)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"[{category}] " if category else "■  "
    _set_font(r1, size, True, color)
    r2 = p.add_run(); r2.text = title
    _set_font(r2, size, True, color)
    return tb


def add_lead(slide, text, size=None):
    """제목 하단 리드메시지: `■ ` 접두 볼드 검정 한 줄 (v4.1, 산출물 실측 반영)."""
    l, t, w, h = LEAD
    return _txt(slide, l, t, w, h, "■ " + text,
                size=size or FONT_PT["lead"], bold=True, color=INK)


def add_bullets(slide, items, l=MARGIN_L, t=BODY_TOP, w=BODY_W,
                line_h=LINE_H, bold_top=True):
    """표준 개조식 불릿 배치 (v4.1). items: [(level, text), ...] → 종료 y 반환.

    수준별 마크·크기: 0 `•` 12pt볼드 / 1 `-` 11pt / 2 `·` 10pt (FONT_PT 스케일 고정).
    □·○ 등 이형 마크, 임의 크기 사용 금지.
    """
    sizes = (FONT_PT["bullet0"], FONT_PT["bullet1"], FONT_PT["bullet2"])
    y = t
    for level, text in items:
        lv = min(level, len(BULLET_MARKS) - 1)
        _txt(slide, l + level * BULLET_INDENT, y, w - level * BULLET_INDENT, line_h,
             f"{BULLET_MARKS[lv]} {text}", size=sizes[lv],
             bold=(level == 0 and bold_top), color=INK)
        y += line_h
    return y


def add_block_header(slide, l, t, w, text):
    """[본장 블록 소제목] 【 제목 】 — 12pt 볼드 검정 (5/15·Link 자료 실측)."""
    return _txt(slide, l, t, w, 0.55, f"【 {text} 】",
                size=FONT_PT["bullet0"], bold=True, color=INK)


def add_hline(slide, l, t, w, color=None, w_pt=0.75):
    """가로 괘선 — 구분 매트릭스(D형) 행 구분·제목 밑줄 룰. (그림자·점선 없음)"""
    return _line(slide, l, t, l + w, t,
                 color=color if color is not None else GRID, w_pt=w_pt)


def add_timeline(slide, l, t, w, milestones):
    """[경과형 E] 가로 타임라인: 굵은 축 + 마일스톤 점, 상단 날짜·하단 라벨.

    milestones: [(날짜, 라벨), ...] → 종료 y 반환. (OSP 경과보고·Coles 협업 실측)
    """
    axis_y = t + 0.55
    ax = _line(slide, l, axis_y, l + w, axis_y, color=INK, w_pt=2.0)
    lnEl = ax.line._get_or_add_ln()                      # 우측 끝 화살촉
    lnEl.append(lnEl.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    step = w / len(milestones)
    for i, (date, label) in enumerate(milestones):
        cx = l + step * (i + 0.5)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(cx - 0.13),
                                     Cm(axis_y - 0.13), Cm(0.26), Cm(0.26))
        dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = INK; dot.line.width = Pt(1.0)
        _txt(slide, cx - step / 2, t, step, 0.45, date, size=FONT_PT["table"],
             bold=True, color=SUB_GRAY, align=PP_ALIGN.CENTER)
        _txt(slide, cx - step / 2, axis_y + 0.20, step, 0.55, label,
             size=FONT_PT["bullet1"], bold=True, color=INK, align=PP_ALIGN.CENTER)
    return axis_y + 0.80


def add_matrix2x2(slide, l, t, w, h, col_labels, row_labels, cells,
                  star=(0, 1), corner="속도＼물량", dim=None):
    """[2×2 포지셔닝 맵] **하나의 3×3 표 객체**로 구현 (v4.1 확정 — 파편화 금지).

    헤더행(0)=물량 축, 헤더열(0)=속도 축, 내부 4셀=서비스 위치. 축 라벨은 헤더 셀에
    녹인다(떠 있는 텍스트박스 금지). star 셀만 음영+네이비 볼드.
    col_labels: [소량, 대용량] / row_labels: [빠름…, 예약…] / cells: [[r0c0,r0c1],[r1c0,r1c1]].
    dim: 의도적 공백 사분면 내부 좌표(예: (1,0) 예약×소량) — 대시 대신 `해당 없음`(옅은 회색·가운데)
         으로 채워 '의도된 공백'을 표시(v4.1 §0.7). 대시 하나만 덩그러니 두지 말 것.
    """
    cells = [list(cells[0]), list(cells[1])]     # 원본 불변 — 복사본에만 기입
    dim_cells = None
    if dim is not None:
        cells[dim[0]][dim[1]] = "해당 없음"
        dim_cells = [(dim[0] + 1, dim[1] + 1)]   # 내부 좌표 → 표 좌표(헤더 +1)
    data = [
        [corner, col_labels[0], col_labels[1]],
        [row_labels[0], cells[0][0], cells[0][1]],
        [row_labels[1], cells[1][0], cells[1][1]],
    ]
    star_cell = (star[0] + 1, star[1] + 1)      # 내부 좌표 → 표 좌표(헤더 +1)
    tbl = add_fin_table(slide, l, t, w, h, data,
                        col_w=[w * 0.30, w * 0.35, w - w * 0.30 - w * 0.35],
                        header_rows=1, header_cols=1, hl_cells=[star_cell],
                        dim_cells=dim_cells)
    # 행 높이: 헤더행 짧게, 내부(사분면) 행 크게 → 맵이 컬럼 세로 공간을 채움
    hdr_h = 0.95
    body_h = (h - hdr_h) / 2
    tbl.rows[0].height = Cm(hdr_h)
    tbl.rows[1].height = Cm(body_h)
    tbl.rows[2].height = Cm(body_h)
    return t + h


def add_vtimeline(slide, l, t, w, h, steps, max_gap=1.55):
    """[세로 타임라인] 세로축 + 단계 점, 각 점 우측 [시점]·[내용] — 좁은 컬럼 세로 공간 채움.

    steps: [(시점, 내용), ...] (좁은 세로 컬럼용 — 좌우 분할 밴드 등, v4.1 실측). → 종료 y 반환.
    max_gap: 노드 간 세로 간격 상한(cm, v4.1 §0.6) — 노드를 억지로 h 전체에 늘리지 않는다.
             간격이 상한을 넘으면 상한으로 조밀하게 두고 남는 공간은 하단 여백으로 비운다.
    """
    n = len(steps)
    dot_x = l + 0.20
    axis_top = t + 0.30
    axis_bot = t + h - 0.30
    gap = (axis_bot - axis_top) / (n - 1) if n > 1 else 0
    if gap > max_gap:                       # 조밀 상한 적용 → 축을 위쪽에 고정, 하단 여백
        gap = max_gap
        axis_bot = axis_top + gap * (n - 1)
    _line(slide, dot_x, axis_top, dot_x, axis_bot, color=INK, w_pt=1.75)
    for i, (when, content) in enumerate(steps):
        cy = axis_top + gap * i
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(dot_x - 0.13),
                                     Cm(cy - 0.13), Cm(0.26), Cm(0.26))
        dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = INK; dot.line.width = Pt(1.0)
        _txt(slide, dot_x + 0.45, cy - 0.45, 1.9, 0.90, when,
             size=FONT_PT["bullet1"], bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _txt(slide, dot_x + 2.45, cy - 0.45, w - dot_x - 2.45 + l, 0.90, content,
             size=FONT_PT["bullet1"], color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return axis_bot + 0.30


def add_htimeline(slide, l, t, w, steps, when_h=0.45, content_h=1.15, size=None):
    """[가로 타임라인] 굵은 축 + 노드 점, 노드 위 [시점](볼드)·아래 [내용] — 전폭 밴드용.

    steps: [(시점, 내용), ...] — 밴드 폭 w 에 노드 균등 분포(세로 늘어짐 없음, v4.1 §4). → 종료 y 반환.
    내용은 여러 줄(\\n) 허용, 노드끼리 겹치지 않게 세그먼트 폭 안에서 가운데 줄바꿈.
    size: 시점·내용 폰트 크기(기본 bullet1=11pt) — 1줄 압축 시 축소해 잘림 방지.
    """
    fs = size or FONT_PT["bullet1"]
    n = len(steps)
    axis_y = t + when_h + 0.12
    ax = _line(slide, l, axis_y, l + w, axis_y, color=INK, w_pt=2.0)
    lnEl = ax.line._get_or_add_ln()                      # 우측 끝 화살촉
    lnEl.append(lnEl.makeelement(qn("a:tailEnd"),
                                 {"type": "triangle", "w": "med", "len": "med"}))
    step = w / n
    for i, (when, content) in enumerate(steps):
        cx = l + step * (i + 0.5)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Cm(cx - 0.14),
                                     Cm(axis_y - 0.14), Cm(0.28), Cm(0.28))
        dot.shadow.inherit = False
        dot.fill.solid(); dot.fill.fore_color.rgb = WHITE
        dot.line.color.rgb = INK; dot.line.width = Pt(1.25)
        _txt(slide, cx - step / 2, t, step, when_h, when, size=fs,
             bold=True, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        _txt(slide, cx - step / 2, axis_y + 0.18, step, content_h, content,
             size=fs, color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    return axis_y + 0.18 + content_h


def add_flow_row(slide, l, t, w, items, h=1.35, gap=0.90):
    """[가로 논리 흐름] 칩(옅은 헤더색) N개를 → 연결로 전폭 배치 — 배경·논거 압축용(v4.1 §5).

    items: [문구, ...] — 좌→우 인과·전개 흐름. 마지막 칩에 결론/공백을 둔다. → 종료 y 반환.
    """
    n = len(items)
    seg_w = (w - gap * (n - 1)) / n
    x = l
    for i, item in enumerate(items):
        _rect(slide, x, t, seg_w, h, fill=TH_PRIMARY)
        _txt(slide, x + 0.25, t, seg_w - 0.50, h, item, size=FONT_PT["bullet1"],
             color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            _txt(slide, x + seg_w, t, gap, h, "→", size=FONT_PT["lead"], bold=True,
                 color=SUB_GRAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += seg_w + gap
    return t + h


def add_specbox(slide, l, t, w, h, rows, header=None, label_w=2.5, hdr_h=0.95):
    """[스펙 패널] 프레임 박스 + (선택)헤더 스트립 + 굵은 라벨:값 행 균등 분배 — 맵과 시각적 쌍둥이.

    rows: [(라벨, 값), ...] — 행을 박스 높이에 균등 분배(헐거운 흩뿌림 방지, 밀도 확보).
    맵(add_matrix2x2)과 동일 t·h 로 나란히 두면 좌우 질량이 맞아 바닥선까지 정렬된다.
    표가 아니라 패널(가로 구분선만, 세로 격자 없음) — '표는 맵 하나' 가드레일과 양립.
    """
    _rect(slide, l, t, w, h, fill=WHITE, line=GRID, line_w=0.75)
    y = t
    if header:
        _rect(slide, l, t, w, hdr_h, fill=TH_PRIMARY)
        _txt(slide, l, t, w, hdr_h, header, size=FONT_PT["table"], bold=True,
             color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        y = t + hdr_h
        _line(slide, l, y, l + w, y, color=GRID, w_pt=0.75)
    n = len(rows)
    pitch = (t + h - y) / n
    for i, (lab, val) in enumerate(rows):
        ry = y + pitch * i
        if i > 0:
            _line(slide, l + 0.15, ry, l + w - 0.15, ry, color=GRID, w_pt=0.5)
        _txt(slide, l + 0.30, ry, label_w, pitch, lab, size=FONT_PT["bullet1"],
             bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
        _txt(slide, l + 0.30 + label_w, ry, w - 0.60 - label_w, pitch, val,
             size=FONT_PT["bullet1"], color=INK, anchor=MSO_ANCHOR.MIDDLE)
    return t + h


def add_frame(slide, l, t, w, h, fill=WHITE, line=None, line_w=0.75):
    """[프레임 박스] 얇은 테두리 박스 — 타임라인·강조 영역을 감싸 밀도·구획감 부여."""
    return _rect(slide, l, t, w, h, fill=fill,
                 line=line if line is not None else GRID, line_w=line_w)


def add_insight_box(slide, l, t, w, h, items, tab="주요 시사점"):
    """[시사점 박스] 회색 테두리 박스 + 좌상단 짙은 탭 라벨 (Coles 협업 실측).

    items: [(level, text), ...] — 내부는 표준 불릿 체계.
    """
    _rect(slide, l, t + 0.28, w, h, fill=WHITE, line=GRID, line_w=1.0)
    _rect(slide, l + 0.35, t, 3.0, 0.55, fill=SUB_GRAY)
    _txt(slide, l + 0.35, t, 3.0, 0.55, tab, size=FONT_PT["table"], bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_bullets(slide, items, l=l + 0.55, t=t + 0.75, w=w - 1.1, line_h=0.55)


def add_col_header(slide, l, t, w, text):
    """[본장 2단 컬럼] 컬럼 헤더 바 — #EAEEF6 채움 + 검정 볼드 가운데 (v4.1 파란 글자 금지)."""
    _rect(slide, l, t, w, COL_HEADER_H, fill=TH_PRIMARY)
    _txt(slide, l, t, w, COL_HEADER_H, text, size=FONT_PT["bullet0"], bold=True,
         color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return t + COL_HEADER_H


def add_conclusion_box(slide, l, t, w, text, h=CONCL_H):
    """[본장 2단 컬럼] 컬럼 하단 결론 박스 — 네이비 채움 + 백색 볼드 가운데 (So-What)."""
    _rect(slide, l, t, w, h, fill=NAVY)
    _txt(slide, l, t, w, h, text, size=FONT_PT["bullet1"], bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def add_vsep(slide, x, t, h):
    """[본장 2단 컬럼] 중앙 세로 구분선 (#BFBFBF 0.75pt 실선, 그림자·점선 없음)."""
    return _line(slide, x, t, x, t + h, color=GRID, w_pt=0.75)


def add_unit(slide, text="(단위 : 억원, %)"):
    """우상단 단위 표기."""
    l, t, w, h = UNIT
    return _txt(slide, l, t, w, h, text, size=FONT_PT["unit"], bold=True,
                align=PP_ALIGN.RIGHT, color=INK)


def add_footnote(slide, text):
    """좌하단 각주 (※ ...)."""
    l, t, w, h = FOOT
    return _txt(slide, l, t, w, h, text, size=FONT_PT["footnote"], bold=False, color=INK)


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])   # 빈 레이아웃


# ─────────────────────────────────────────────────────────────
# 4. 원형(Archetype) 빌더
# ─────────────────────────────────────────────────────────────
def new_deck():
    """A4 가로 빈 프레젠테이션."""
    prs = Presentation()
    prs.slide_width = Emu(9906000)
    prs.slide_height = Emu(6858000)
    return prs


def add_cover(prs, title, org="온라인사업단", date="2026. 07. 06", tag="내부 토의용"):
    """[표지] 상단 로고존 · 태그 · 대제목(밑줄) · 하단 조직/일자."""
    s = _blank(prs)
    _txt(s, 2.0, 6.4, 10.0, 0.7, tag, size=13, bold=True, color=INK)
    # 대제목 (가운데, 밑줄)
    tb = _txt(s, 2.0, 7.3, 23.5, 2.0, title, size=40, bold=True, color=INK,
              align=PP_ALIGN.CENTER)
    # 밑줄
    ln = s.shapes.add_connector(2, Cm(2.2), Cm(9.4), Cm(25.3), Cm(9.4))
    ln.line.color.rgb = INK; ln.line.width = Pt(1.5)
    _txt(s, 2.0, 14.3, 23.5, 0.7, org, size=14, bold=True, color=INK,
         align=PP_ALIGN.CENTER)
    _txt(s, 2.0, 15.0, 23.5, 0.7, date, size=13, bold=True, color=INK,
         align=PP_ALIGN.CENTER)
    return s


def add_divider(prs, num, title, items):
    """[챕터 간지] 회색 풀블리드 + 대제목 + ①②③ 하위항목."""
    s = _blank(prs)
    _rect(s, 0.0, 3.65, PAGE_W_CM, 13.9, fill=DIVIDER_BG)
    _txt(s, 3.5, 5.6, 20.5, 1.2, f"{num}. {title}", size=28, bold=True,
         color=INK, align=PP_ALIGN.CENTER)
    circ = "①②③④⑤⑥⑦⑧⑨"
    y = 8.0
    for i, it in enumerate(items):
        _txt(s, 9.0, y, 15.0, 0.7, f"{circ[i]} {it}", size=15, bold=True, color=INK)
        y += 0.85
    return s


def add_content(prs, category, title, tier="P",
                lead=None, unit=None, footnote=None):
    """[본문 P/F계열] 크롬만 세팅한 슬라이드 반환 → 본문 자유 배치.

    lead: 제목 하단 `■ ` 리드메시지 한 줄 (선택). 사용 시 본문은 BODY_TOP_LEAD(4.10)부터.
    페이지번호 크롬은 v4.1 에서 제거.
    """
    s = _blank(prs)
    add_header(s, category, title, tier)
    if lead:
        add_lead(s, lead)
    if unit:
        add_unit(s, unit)
    if footnote:
        add_footnote(s, footnote)
    return s


def add_closing(prs, company="롯데마트 · 롯데슈퍼"):
    """[종료] No.1 GROCERY MARKET 마감 장표(간이)."""
    s = _blank(prs)
    _txt(s, 2.0, 2.0, 20.0, 2.2, "No.1", size=72, bold=True, color=CRIMSON)
    _txt(s, 2.0, 4.6, 24.0, 1.6, "GROCERY MARKET", size=44, bold=True, color=CRIMSON)
    _txt(s, 2.05, 7.0, 20.0, 0.8, "Discover a joyful food life", size=18,
         bold=False, color=RGBColor(0x80, 0x80, 0x80))
    cols = [("RE:TARGET", "경영목표"), ("RE:DESIGN", "고객경험"), ("RE:FOCUS", "일하는 방식")]
    xs = [1.5, 10.0, 18.5]
    for (en, ko), x in zip(cols, xs):
        ln = s.shapes.add_connector(2, Cm(x), Cm(13.0), Cm(x + 6.0), Cm(13.0))
        ln.line.color.rgb = INK; ln.line.width = Pt(1.0)
        tb = _txt(s, x, 13.2, 7.0, 0.7, "", size=15, bold=True, color=NAVY)
        p = tb.text_frame.paragraphs[0]
        r1 = p.add_run(); r1.text = en + "  "; _set_font(r1, 15, True, NAVY)
        r2 = p.add_run(); r2.text = ko; _set_font(r2, 10, False, INK)
    return s


# ─────────────────────────────────────────────────────────────
# 5. 재무형 표 (EAEEF6 헤더 · 전 셀 가운데 · 재무 수치열만 우측)
# ─────────────────────────────────────────────────────────────
_TBL_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}


def add_fin_table(slide, l, t, w, h, data, col_w=None,
                  header_rows=1, header_fill=TH_PRIMARY,
                  col_align=None, font_size=None,
                  merges=None, hl_rows=None, hl_cols=None, hl_cells=None,
                  dim_cells=None, header_cols=0, bold_cols=(0,), bold_rows=None):
    """
    data: 2D list [ [row0col0, ...], ... ]  (문자열, None=위 셀과 병합 예정 자리)
    header_rows: 상단 헤더 행 수 (헤더 채움 + 검정 볼드 + 가운데)
    merges: [(row_start, row_end, col), ...] — 동일 값·그룹 구분의 세로 병합 (v4.1)
    hl_rows: 강조 행 인덱스 — 연블루(HL_FILL) 채움 + 네이비(HL_TEXT) 볼드 (신규·핵심 행)
    hl_cols: 강조 열 인덱스 — 신·구 비교표에서 신규 서비스 열 음영 (헤더 포함, v4.1)
    bold_rows: 볼드만 적용할 행 인덱스 (음영 없이 텍스트 볼드 — 핵심 대비 행 강조 역할 분리)
    dim_cells: 의도된 공백 셀 [(i,j),...] — 옅은 회색(SUB_GRAY)·볼드 해제 (v4.1 §0.7)
    bold_cols: 볼드 처리 열 (기본 1열; 서비스명 등 구분성 열 추가 지정 가능)

    정렬 표준 (v4.1, 정기협의체 본장 실측): 헤더·본문 전 셀 **가운데** 고정.
    재무 수치열만 col_align 으로 우측 지정 — 예: col_align=["c","r","r"].
    셀별 임의 정렬 금지.
    """
    rows, cols = len(data), len(data[0])
    gtbl = slide.shapes.add_table(rows, cols, Cm(l), Cm(t), Cm(w), Cm(h))
    tbl = gtbl.table
    # 기본 스타일 제거(밴딩 off)
    tbl.first_row = False; tbl.horz_banding = False
    if col_w:
        for j, cw in enumerate(col_w):
            tbl.columns[j].width = Cm(cw)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Cm(0.1); cell.margin_right = Cm(0.1)
            cell.margin_top = 0; cell.margin_bottom = 0
            _cell_border(cell)                       # 격자 #BFBFBF 0.75pt (v4.1 코드 강제)
            is_head = i < header_rows or j < header_cols
            is_hl = ((hl_rows and i in hl_rows) or (hl_cols and j in hl_cols)
                     or (hl_cells and (i, j) in hl_cells))
            if is_hl:
                cell.fill.solid(); cell.fill.fore_color.rgb = HL_FILL
            elif is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame; tf.word_wrap = True
            if is_head:
                align = PP_ALIGN.CENTER
            else:
                align = _TBL_ALIGN[col_align[j] if col_align else "c"]
            is_dim = bool(dim_cells and (i, j) in dim_cells)
            # 다중행 셀도 전 행 동일 정렬 — 줄마다 문단 분리 (None = 병합 자리)
            for k, line in enumerate(str(val).split("\n") if val is not None else []):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = align
                _emit_runs(p, line, font_size or FONT_PT["table"],
                           (not is_dim) and (is_head or is_hl or j in bold_cols
                                             or (bold_rows and i in bold_rows)),
                           SUB_GRAY if is_dim else (HL_TEXT if is_hl else INK))
    if merges:
        for r0, r1, jc in merges:
            tbl.cell(r0, jc).merge(tbl.cell(r1, jc))
    return tbl
