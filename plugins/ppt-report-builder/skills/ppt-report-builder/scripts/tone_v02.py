# -*- coding: utf-8 -*-
"""tone_v02 — 임원 보고 'v0.2' 톤 헬퍼(네이비 번호박스 헤더·빨강 이탤릭 강조·차트).

zetta_ppt_standard 위에 얹는 별도 디자인 톤. import 예:
    from zetta_ppt_standard import *
    from tone_v02 import *
    prs = new_deck(); s = new_slide(prs); reset_pageno()
    hdr(s, 1, "섹션명_슬라이드 주제"); lead(s, "리딩 한 줄"); pageno(s)

색: NAVY2 #1F3864 / RED #C00000 / GRAY #808080. 페이지 A4 가로(new_deck).
"""
from zetta_ppt_standard import (
    _txt, _rect, _line, _set_font, new_deck,
    RGBColor, Cm, Pt, qn, WHITE, INK,
    PP_ALIGN, MSO_ANCHOR, MSO_SHAPE,
)

NAVY2 = RGBColor(0x1F, 0x38, 0x64)
RED = RGBColor(0xC0, 0x00, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)
REDBG = RGBColor(0xFD, 0xE9, 0xE9)
STRIP = RGBColor(0xEE, 0xF1, 0xF7)
X0, XR, CW = 1.2, 26.32, 25.12
_pg = [0]


def reset_pageno():
    _pg[0] = 0


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def hdr(s, secno, title):
    """좌상단 네이비 번호박스 + 굵은 제목 + 네이비 실선. 제목은 '섹션명_주제' 형식 권장."""
    _rect(s, X0, 0.85, 1.05, 1.05, fill=NAVY2)
    _txt(s, X0, 0.85, 1.05, 1.05, str(secno), size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, 2.5, 0.85, CW - 1.3, 1.05, title, size=19, bold=True, color=NAVY2,
         anchor=MSO_ANCHOR.MIDDLE)
    _line(s, X0, 2.02, XR, 2.02, color=NAVY2, w_pt=2.25)


def lead(s, text):
    """제목 아래 리딩 1줄(개조식 명사 종결)."""
    _txt(s, X0, 2.22, CW, 0.85, text, size=15, bold=True, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)


def redconcl(s, text, y=16.55):
    """하단 중앙 빨강 이탤릭 + 따옴표 강조 문구."""
    tb = _txt(s, X0, y, CW, 0.8, f"“ {text} ”", size=15, bold=True,
              color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.italic = True


def foot(s, text):
    """좌하단 각주(8~9pt 회색, ※/Source 시작)."""
    _txt(s, X0, 17.9, 21.0, 0.7, text, size=9, color=GRAY)


def pageno(s):
    """우하단 페이지번호(모듈 카운터 자동 증가; new_slide 마다 pageno 호출)."""
    _pg[0] += 1
    _txt(s, 24.3, 17.9, 2.0, 0.6, str(_pg[0]), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def box(s, x, y, w, h, title, body, tfill=NAVY2, bfill=WHITE, border=NAVY2):
    """제목 스트립(네이비/빨강) + 본문 카드 박스."""
    _rect(s, x, y, w, h, fill=bfill, line=border, line_w=1.0)
    _rect(s, x, y, w, 0.85, fill=tfill)
    _txt(s, x, y, w, 0.85, title, size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, x + 0.35, y + 1.0, w - 0.7, h - 1.2, body, size=11, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)


def tbl(s, x, y, col_w, data, aligns, hl_row=None, header_h=0.95, body_h=1.1,
        head_fill=NAVY2, fs=11):
    """네이비 헤더·짝수행 음영 표. hl_row=강조 행(옅은 빨강). data[0]=헤더."""
    from zetta_ppt_standard import _cell_border
    rows, cols = len(data), len(data[0])
    total = header_h + (rows - 1) * body_h
    t = s.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(sum(col_w)), Cm(total)).table
    t.first_row = False; t.horz_banding = False
    for j, cw in enumerate(col_w):
        t.columns[j].width = Cm(cw)
    t.rows[0].height = Cm(header_h)
    for i in range(1, rows):
        t.rows[i].height = Cm(body_h)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Cm(0.2); c.margin_right = Cm(0.12)
            c.margin_top = 0; c.margin_bottom = 0
            _cell_border(c)
            f = head_fill if i == 0 else (REDBG if hl_row == i else (STRIP if i % 2 == 0 else WHITE))
            c.fill.solid(); c.fill.fore_color.rgb = f
            tf = c.text_frame; tf.word_wrap = True
            al = PP_ALIGN.CENTER if i == 0 else aligns[j]
            for k, ln in enumerate(str(val).split("\n")):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = al
                r = p.add_run(); r.text = ln
                _set_font(r, 12 if i == 0 else fs, i == 0, WHITE if i == 0 else INK)
    return y + total


def arrow(s, x1, y1, x2, y2, color=RED, w=3.0):
    """화살표 커넥터(그림자 없음)."""
    ln = s.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    ln.shadow.inherit = False; ln.line.color.rgb = color; ln.line.width = Pt(w)
    el = ln.line._get_or_add_ln()
    el.append(el.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def flownode(s, x, y, w, h, text, fill=NAVY2, tc=WHITE, size=12):
    """플로우 도식 노드(라운드 사각형·가운데 정렬)."""
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.shadow.inherit = False; sp.adjustments[0] = 0.06
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = fill
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Cm(0.08))
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, size, True, tc)
    return sp


def column_chart(s, x, y, w, h, categories, values, per_point_colors=None,
                 number_format='0', label_size=11):
    """세로 막대 차트(범례·제목 없음, 데이터 라벨 볼드). per_point_colors=[RGB,...] 선택.
    〔조사〕 미확정 축은 categories=['__',...] 로 비워 둔다(임의 수치 금지)."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
    cd = CategoryChartData(); cd.categories = categories; cd.add_series("s", values)
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = False
    pl = ch.plots[0]; pl.has_data_labels = True
    dl = pl.data_labels; dl.number_format = number_format; dl.number_format_is_linked = False
    dl.font.size = Pt(label_size); dl.font.bold = True
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    if per_point_colors:
        for i, pt in enumerate(pl.series[0].points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = per_point_colors[i]
    return gf


def bar_chart(s, x, y, w, h, categories, values, per_point_colors=None,
              number_format='+0.0"%"', label_size=11):
    """가로 막대 차트(범례·제목 없음). categories 는 아래→위 순으로 그려진다."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    cd = CategoryChartData(); cd.categories = categories; cd.add_series("s", values)
    gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Cm(x), Cm(y), Cm(w), Cm(h), cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = False
    pl = ch.plots[0]; pl.has_data_labels = True
    dl = pl.data_labels; dl.number_format = number_format; dl.number_format_is_linked = False
    dl.font.size = Pt(label_size); dl.font.bold = True
    if per_point_colors:
        for i, pt in enumerate(pl.series[0].points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = per_point_colors[i]
    return gf
