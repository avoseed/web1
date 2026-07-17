# -*- coding: utf-8 -*-
"""롯데마트 캄보디아 글로벌 확장 전략 v0.3 — 12장(표지·목차 포함), v0.2 디자인 톤.

톤: 좌상단 네이비 번호 박스 + 굵은 제목 + 네이비 실선 / 리딩 1줄(명사 종결) /
    하단 중앙 빨강 이탤릭 강조 / 네이비 #1F3864·빨강 #C00000 / 우하단 페이지번호.
가드레일: 〔조사〕 칸은 공란/__ 유지(임의 수치 금지) · 데이터 슬라이드 Source 각주 ·
          지도 이미지 금지(자체 도식) · 오탈자 수정(신선식품·운송·프랜차이즈).
"""
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "zetta-ppt-standard"))
from zetta_ppt_standard import *
from zetta_ppt_standard import _txt, _rect, _line, _set_font
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

NAVY2 = RGBColor(0x1F, 0x38, 0x64)
RED = RGBColor(0xC0, 0x00, 0x00)
GRAY = RGBColor(0x80, 0x80, 0x80)
REDBG = RGBColor(0xFD, 0xE9, 0xE9)
STRIP = RGBColor(0xEE, 0xF1, 0xF7)
X0, XR, CW = 1.2, 26.32, 25.12
_pg = [0]


def newslide():
    return prs.slides.add_slide(prs.slide_layouts[6])


def hdr(s, secno, title):
    _rect(s, X0, 0.85, 1.05, 1.05, fill=NAVY2)
    _txt(s, X0, 0.85, 1.05, 1.05, str(secno), size=20, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, 2.5, 0.85, CW - 1.3, 1.05, title, size=19, bold=True, color=NAVY2,
         anchor=MSO_ANCHOR.MIDDLE)
    _line(s, X0, 2.02, XR, 2.02, color=NAVY2, w_pt=2.25)


def lead(s, text):
    _txt(s, X0, 2.22, CW, 0.85, text, size=15, bold=True, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)


def redconcl(s, text, y=16.55):
    tb = _txt(s, X0, y, CW, 0.8, f"“ {text} ”", size=15, bold=True,
              color=RED, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            r.font.italic = True


def foot(s, text):
    _txt(s, X0, 17.9, 21.0, 0.7, text, size=9, color=GRAY)


def pageno(s):
    _pg[0] += 1
    _txt(s, 24.3, 17.9, 2.0, 0.6, str(_pg[0]), size=9, color=GRAY, align=PP_ALIGN.RIGHT)


def box(s, x, y, w, h, title, body, tfill=NAVY2, bfill=WHITE, border=NAVY2):
    _rect(s, x, y, w, h, fill=bfill, line=border, line_w=1.0)
    _rect(s, x, y, w, 0.85, fill=tfill)
    _txt(s, x, y, w, 0.85, title, size=12, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, x + 0.35, y + 1.0, w - 0.7, h - 1.2, body, size=11, color=INK,
         anchor=MSO_ANCHOR.MIDDLE)


def tbl(s, x, y, col_w, data, aligns, hl_row=None, header_h=0.95, body_h=1.1,
        head_fill=NAVY2, fs=11):
    from zetta_ppt_standard import _cell_border
    rows, cols = len(data), len(data[0])
    total = header_h + (rows - 1) * body_h
    t = s.shapes.add_table(rows, cols, Cm(x), Cm(y), Cm(sum(col_w)), Cm(total)).table
    t.first_row = False; t.horz_banding = False
    for j, cw in enumerate(col_w):
        t.columns[j].width = Cm(cw)
    t.rows[0].height = Cm(header_h)
    for i in range(1, rows):
        t.rows[i].height = Cm(body_h)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.vertical_anchor = MSO_ANCHOR.MIDDLE
            c.margin_left = Cm(0.2); c.margin_right = Cm(0.12)
            c.margin_top = 0; c.margin_bottom = 0
            _cell_border(c)
            f = head_fill if i == 0 else (REDBG if hl_row == i else (STRIP if i % 2 == 0 else WHITE))
            c.fill.solid(); c.fill.fore_color.rgb = f
            tf = c.text_frame; tf.word_wrap = True
            al = PP_ALIGN.CENTER if i == 0 else aligns[j]
            for k, ln in enumerate(str(val).split("\n")):
                p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
                p.alignment = al
                r = p.add_run(); r.text = ln
                _set_font(r, 12 if i == 0 else fs, i == 0, WHITE if i == 0 else INK)
    return y + total


def arrow(s, x1, y1, x2, y2, color=RED, w=3.0):
    ln = s.shapes.add_connector(2, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    ln.shadow.inherit = False; ln.line.color.rgb = color; ln.line.width = Pt(w)
    el = ln.line._get_or_add_ln()
    el.append(el.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def flownode(s, x, y, w, h, text, fill=NAVY2, tc=WHITE, size=12):
    sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    sp.shadow.inherit = False; sp.adjustments[0] = 0.06
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = fill
    tf = sp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
        setattr(tf, m, Cm(0.08))
    for i, ln in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, size, True, tc)
    return sp


prs = new_deck()

# ══════════ 표지 ══════════
s = newslide()
_txt(s, 1.5, 0.9, 12.0, 0.6, "Lifetime Value Creator", size=11, bold=True, color=NAVY2)
_txt(s, 1.5, 7.0, 24.5, 1.6, "롯데마트 캄보디아 글로벌 확장 전략", size=30, bold=True,
     color=NAVY2, align=PP_ALIGN.CENTER)
_line(s, 8.0, 8.9, 19.5, 8.9, color=RED, w_pt=2.0)
_txt(s, 1.5, 9.2, 24.5, 1.0, "베트남 검증 역량 기반 Asset-Light 진입 방안", size=16,
     bold=True, color=INK, align=PP_ALIGN.CENTER)
_txt(s, 1.5, 15.5, 24.5, 0.8, "2026. 07. __", size=13, color=GRAY, align=PP_ALIGN.CENTER)
_txt(s, 22.5, 17.9, 3.8, 0.6, "[ 롯데 심볼 ]", size=9, color=GRAY, align=PP_ALIGN.RIGHT)

# ══════════ 목차 ══════════
s = newslide()
hdr(s, "C", "목차  |  Contents")
items = ["롯데마트 베트남 현황", "글로벌 확장 전략 배경", "핵심역량 강화", "향후 계획 · 제언"]
yy = 4.2
for i, it in enumerate(items, 1):
    _rect(s, 4.0, yy, 1.0, 1.0, fill=NAVY2)
    _txt(s, 4.0, yy, 1.0, 1.0, str(i), size=18, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, 5.4, yy, 17.0, 1.0, it, size=18, bold=True, color=NAVY2, anchor=MSO_ANCHOR.MIDDLE)
    yy += 2.6
pageno(s)

# ══════════ 1-1 베트남 현황_실적 및 M/S ══════════
s = newslide()
hdr(s, 1, "롯데마트 베트남 현황_실적 및 M/S")
lead(s, "M/S 성장 중이나 경쟁사 대비 저조 — 현 구조만으로는 M/S 잠식 우려")
# 좌: 영업이익 추이 컬럼 차트 (▲11 → 180 → 231 → 300 → 371)
_txt(s, X0, 3.2, 12.0, 0.6, "베트남 법인 영업이익 추이   (단위 : 억 원)",
     size=12, bold=True, color=NAVY2)
cd = CategoryChartData()
cd.categories = ["__", "__", "__", "__", "__"]   # 〔조사: 연도 확인 후 기입〕
cd.add_series("영업이익", (-11, 180, 231, 300, 371))
gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Cm(X0), Cm(3.9), Cm(12.0), Cm(9.5), cd)
ch = gf.chart; ch.has_legend = False; ch.has_title = False
pl = ch.plots[0]; pl.has_data_labels = True
pl.data_labels.number_format = '0;▲0'; pl.data_labels.number_format_is_linked = False
pl.data_labels.font.size = Pt(11); pl.data_labels.font.bold = True
pl.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
for i, pt in enumerate(pl.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RED if i == 0 else NAVY2
_txt(s, X0, 13.5, 12.0, 0.5, "〔조사 : 각 수치의 해당 연도 확인 후 기입〕  ▲ = 영업손실",
     size=9, color=GRAY)
# 우: 경쟁사 매출성장률 가로 막대
_txt(s, 14.2, 3.2, 12.0, 0.6, "경쟁사 매출성장률 비교", size=12, bold=True, color=NAVY2)
cd2 = CategoryChartData()
cd2.categories = ["Lotte", "Aeon", "Masan(WinMart)", "BHX(Mobile World)", "Central (GO!/mini go!)"]
cd2.add_series("성장률", (2.2, 13.9, 13.9, 14.1, 56.0))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Cm(14.2), Cm(3.9), Cm(12.1), Cm(9.5), cd2)
ch2 = gf2.chart; ch2.has_legend = False; ch2.has_title = False
pl2 = ch2.plots[0]; pl2.has_data_labels = True
pl2.data_labels.number_format = '+0.0"%"'; pl2.data_labels.number_format_is_linked = False
pl2.data_labels.font.size = Pt(11); pl2.data_labels.font.bold = True
for i, pt in enumerate(pl2.series[0].points):
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = RED if i == 0 else NAVY2
foot(s, "※ 기준 : 〔조사 : 기준연도〕 매출성장률   ·   Source : 각사 공시")
pageno(s)

# ══════════ 1-2 베트남 현황_경쟁 현황 ══════════
s = newslide()
hdr(s, 1, "롯데마트 베트남 현황_경쟁 현황")
lead(s, "기존 Format 경쟁뿐 아니라 이커머스·편의점·SSM의 성장가능성 高")
tbl(s, X0, 3.5, [4.2, 8.0], [
    ["구분", "주요 경쟁사"],
    ["하이퍼마켓", "GO! · Tops market"],
    ["슈퍼마켓", "AEON · MaxValu"],
], aligns=[PP_ALIGN.CENTER, PP_ALIGN.LEFT], header_h=1.0, body_h=2.1)
_txt(s, X0, 3.05, 12.0, 0.5, "베트남 주요 경쟁사 (Format 별)", size=11, bold=True, color=NAVY2)
_txt(s, 14.2, 3.05, 12.0, 0.5, "신채널 성장 전망", size=11, bold=True, color=NAVY2)
tbl(s, 14.2, 3.5, [4.2, 4.0, 3.9], [
    ["채널", "現 시장 규모", "'26~'30 성장률"],
    ["이커머스", "48조", "15~20%"],
    ["편의점", "0.6조", "11~13%"],
    ["SSM", "9조", "13~15%"],
], aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.CENTER], header_h=1.0, body_h=1.55)
redconcl(s, "내부 효율 → 외부 확장으로 전환 필요")
foot(s, "Source : 각사 공시 · 유통업계 자료")
pageno(s)

# ══════════ 2-1 확장 배경_지역 선정 ══════════
s = newslide()
hdr(s, 2, "글로벌 확장 전략 배경_외부 확장 지역 선정")
lead(s, "캄보디아는 '저위험 + 고성장 + 베트남 연계' 유일 시장")
_txt(s, X0, 3.05, 12.0, 0.5, "메콩 3국 비교 매트릭스", size=11, bold=True, color=NAVY2)
tbl(s, X0, 3.5, [2.6, 2.0, 2.0, 1.9, 2.5, 1.8], [
    ["국가", "시장성", "리스크", "경쟁", "베트남 시너지", "종합"],
    ["캄보디아", "高", "低", "低", "高", "★★★★"],
    ["라오스", "低", "中", "低", "中", "★★"],
    ["미얀마", "高", "高", "中", "低", "★"],
], aligns=[PP_ALIGN.CENTER] * 6, hl_row=1, header_h=1.1, body_h=1.9)
_txt(s, X0, 10.4, 12.8, 1.6,
     "※ 평가 기준 — 시장성 : 소매시장 성장률·인구 / 리스크 : 정치·외환 / "
     "시너지 : 베트남 육로 접근성", size=9, color=GRAY, anchor=MSO_ANCHOR.TOP)
_txt(s, 14.2, 3.05, 12.0, 0.5, "캄보디아 스냅샷", size=11, bold=True, color=NAVY2)
snaps = ["인구 1,689만('23) · 프놈펜 228만",
         "경제성장 1998년 이후 장기평균 7%",
         "달러 통용 경제 → 환위험 완화",
         "목바이·바벳 육로 국경 → 베트남 남부 직결"]
sy = 3.6
for i, sn in enumerate(snaps, 1):
    _rect(s, 14.2, sy, 12.1, 1.7, fill=STRIP, line=NAVY2, line_w=0.75)
    _rect(s, 14.2, sy, 1.3, 1.7, fill=NAVY2)
    _txt(s, 14.2, sy, 1.3, 1.7, str(i), size=16, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, 15.8, sy, 10.3, 1.7, sn, size=12, bold=True, color=INK, anchor=MSO_ANCHOR.MIDDLE)
    sy += 1.95
foot(s, "Source : 캄보디아 신규진출 검토 보고('24.03), KOTRA")
pageno(s)

# ══════════ 2-2 확장 배경_캄보디아 유통 구조 ══════════
s = newslide()
hdr(s, 2, "글로벌 확장 전략 배경_캄보디아 유통 구조")
lead(s, "현대 유통 침투율 저조 → 시장 잠재력 高, 통합 공급 사업자 부재")
# 피라미드 3층
for w, yy, txt, fill, tc in [
        (6.5, 3.6, "E-Commerce\n(초기 단계)", RGBColor(0x9D, 0xB0, 0xCE), INK),
        (9.8, 6.3, "Modern Trade\n슈퍼 545 · 편의점 118 · 하이퍼 5 ('22)", RGBColor(0x4E, 0x6A, 0x9E), WHITE),
        (13.0, 9.0, "Traditional Market\n소규모 식료품점 109,222개 (95.5%)", NAVY2, WHITE)]:
    cxp = 7.6
    sp = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID, Cm(cxp - w / 2), Cm(yy), Cm(w), Cm(2.5))
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = fill; sp.line.color.rgb = fill
    tf = sp.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    for k, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if k == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = ln; _set_font(r, 12 if k == 0 else 10, k == 0, tc)
_txt(s, 0.6, 6.3, 6.6, 2.5, "← 편의점 CAGR 25.2%\n     하이퍼 13.6%",
     size=10, bold=True, color=RED, anchor=MSO_ANCHOR.MIDDLE)
box(s, 15.5, 3.5, 10.8, 3.0, "과반수 수입 의존 구조",
    "'22년 식품 수입 : 태국 23.3% · 인니 13.4% · 베트남 11.1%")
box(s, 15.5, 6.7, 10.8, 3.0, "현대유통 = 프놈펜 집중·외국계 리드",
    "Lucky 85점 · Circle K 39 · AEON 30 · Big C 20")
box(s, 15.5, 9.9, 10.8, 3.0, "수입 → 유통 → 소매 파편화",
    "통합 공급 역량 보유 사업자 부재", tfill=RED)
foot(s, "Source : 캄보디아 신규진출 검토 보고('24.03), KOTRA 유통구조 조사('22)")
pageno(s)

# ══════════ 3-1 핵심역량_시장 기회 요소 ══════════
s = newslide()
hdr(s, 3, "핵심역량 기반 공략_시장 기회 요소")
lead(s, "롯데마트의 '유통 플랫폼' 역할로서 시장 문제 해결")
tbl(s, X0, 3.4, [3.7, 3.0, CW - 6.7], [
    ["시장 문제", "롯데 강점", "세부 문안"],
    ["물류 비효율", "WMS",
     "수출형 창고관리·가시성 솔루션 강화 — 베트남 물류센터↔캄보디아 매장 재고 Data 실시간 연동 / "
     "재고정확도 글로벌 99.8% 이상 / 프랜차이즈 매장 결품율 2.0% 이하 안전재고 최적화"],
    ["신선품질 불균일", "SCM",
     "Cold Chain·육로 수송 → 최적 동선·적재율 향상(Back Haul) / 다중 온도대 차량·TMS 고도화 / "
     "프놈펜 보관·통관·육로 운송 리드타임 단축 — 캄보디아 물류 106위(도로 99·항만 81) → SCM 자체가 진입장벽"],
    ["상품 부족", "Sourcing Power",
     "PB 개발·글로벌 직소싱 — PB 구성비 現 15% → 30% / K-Food 소싱 허브(롯데웰푸드·크라운·동원 F&B·중소 제조사 동반 진출)"],
    ["운영 미숙", "운영 노하우",
     "베트남 매장 운영 Manual 전수 — 신선 품질·위생·상품화 / Deli 특화(Yorihada) Know-how·Recipe / POG 진열 원칙"],
], aligns=[PP_ALIGN.CENTER, PP_ALIGN.CENTER, PP_ALIGN.LEFT], header_h=0.95, body_h=2.6, fs=10)
redconcl(s, "시장 문제 → 롯데의 기회로 전환")
pageno(s)

# ══════════ 3-2 핵심역량_크로스보더 공급망 ══════════
s = newslide()
hdr(s, 3, "핵심역량 기반 공략_크로스보더 공급망")
lead(s, "베트남 남부 인프라를 캄보디아로 연장하는 국경 연계 공급망")
# 좌: 노드 3개 수평 플로우 (지도 금지 — 자체 도형)
nodes = [("호치민\n(소싱·물류 허브)", 1.2), ("목바이 국경\n(통관)", 5.9), ("프놈펜\n(매장·파트너)", 10.6)]
for txt, x in nodes:
    flownode(s, x, 6.6, 3.6, 2.4, txt, size=11)
for x in (4.85, 9.55):
    arrow(s, x, 7.8, x + 1.0, 7.8)
_txt(s, 4.7, 5.6, 1.6, 0.9, "__시간", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
_txt(s, 9.4, 5.6, 1.6, 0.9, "__시간", size=11, bold=True, color=RED, align=PP_ALIGN.CENTER)
_txt(s, 1.2, 9.2, 13.0, 1.0, "❄ 콜드체인 유지 구간   ·   리드타임 〔조사 : D-10 물류 재설계 시 확정〕",
     size=10, bold=True, color=NAVY2)
box(s, 15.0, 3.4, 11.3, 3.0, "공급 거점",
    "베트남 남부 소싱/물류 인프라 기 구축 → 신규 투자 최소")
box(s, 15.0, 6.6, 11.3, 3.0, "육로 연계",
    "국경 통관 → 프놈펜 직결, 콜드체인 유지 배송")
box(s, 15.0, 9.8, 11.3, 3.0, "역량 재정의",
    "SCM → Cross-border Supply Network / PB → Profit Engine", tfill=RED)
pageno(s)

# ══════════ 4-1 향후 계획_수익 모델 ══════════
s = newslide()
hdr(s, 4, "향후 계획_수익 모델")
lead(s, "출점 없이 B2B·PB·로열티·물류로 수익 — 저투자·고마진")
flownode(s, 1.2, 6.0, 5.4, 3.2, "베트남 남부\n소싱·물류\n(국경 연계)", size=12)
arrow(s, 6.7, 7.6, 8.3, 7.6)
flownode(s, 8.4, 6.0, 5.2, 3.2, "B2B 도매\n(초기 현금흐름)", fill=RED, size=12)
rev = [("① 상품공급 마진", 3.6), ("② PB 수익 (고마진)", 6.1), ("③ 마스터 프랜차이즈 로열티", 8.6),
       ("④ 크로스보더 물류 수익", 11.1)]
for label, yy in rev:
    arrow(s, 13.7, 7.6, 15.4, yy + 0.9)
    flownode(s, 15.5, yy, 10.8, 1.8, label, size=12)
_rect(s, X0, 14.4, CW, 1.2, fill=NAVY2)
_txt(s, X0, 14.4, CW, 1.2, "Low Investment  ·  High Margin  ·  Fast Expansion",
     size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pageno(s)

# ══════════ 4-2 향후 계획_단계별 실행 전략 ══════════
s = newslide()
hdr(s, 4, "향후 계획_단계별 실행 전략")
lead(s, "Low risk → High return 순 단계적 진입, 정량 Gate 충족 시 단계 전환")
tbl(s, X0, 3.4, [5.0, 9.0, CW - 14.0], [
    ["단계", "내용", "전환 Gate  〔조사 : D-8 시장조사 후 확정〕"],
    ["Step 1. B2B 상품 공급", "저위험 진입 · 현금흐름 확보", "월 공급액 __억 · 활성 파트너 __개 · 공급마진 __%"],
    ["Step 2. Franchise", "마스터 프랜차이즈 무자본 확장", "가맹 점포 __개 · 점당 월매출 __ · 로열티 안정화 __개월"],
    ["Step 3. Platform", "유통 + 물류 생태계 완성", "— (장기 비전)"],
], aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.LEFT], header_h=1.0, body_h=1.9, fs=11)
box(s, X0, 10.4, CW, 2.5, "파트너 후보군 (Step 2)",
    "Chip Mong(재계 2위) · Royal Group(재계 1위) · 부영그룹\n"
    "※ 부영 : 프놈펜 CBD 임점 제안('24.02) 실재 — 1호 거점 옵션(선택), 본류는 B2B 유지",
    tfill=NAVY2)
foot(s, "※ Gate 수치는 D-8 캄보디아 수출 시장조사 결과로 확정\n"
        "※ 캄보디아 프랜차이즈 법제 미비(제정 준비 중)·외국인 토지 소유 불가(최대 50년 사용권) "
        "→ Step 2 계약 구조는 법무 검토 병행")
pageno(s)

# ══════════ 4-3 향후 계획_1년 로드맵 ══════════
s = newslide()
hdr(s, 4, "향후 계획_1년 로드맵")
lead(s, "캄보디아 진입 1년 내 사업 개시 가능")
steps = [("D-12", "Project Kick-off"), ("D-11", "조직 설계(상품·물류·OP)"),
         ("D-10", "물류 인프라 재설계"), ("D-8", "캄보디아 수출 시장 조사\n→ Gate 수치 확정"),
         ("D-3", "캄보디아 상품 박람회"), ("D-day", "B2B Wholesale 개시")]
axis_x = 2.2
ytop, ybot = 3.7, 13.3
_line(s, axis_x, ytop, axis_x, ybot, color=NAVY2, w_pt=2.0)
gap = (ybot - ytop) / (len(steps) - 1)
for i, (d, t) in enumerate(steps):
    cy = ytop + gap * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Cm(axis_x - 0.16), Cm(cy - 0.16), Cm(0.32), Cm(0.32))
    dot.shadow.inherit = False; dot.fill.solid()
    dot.fill.fore_color.rgb = RED if i == len(steps) - 1 else NAVY2; dot.line.fill.background()
    _txt(s, axis_x + 0.4, cy - 0.4, 2.0, 0.8, d, size=12, bold=True,
         color=NAVY2, anchor=MSO_ANCHOR.MIDDLE)
    _txt(s, axis_x + 2.6, cy - 0.5, 8.5, 1.0, t, size=11, color=INK, anchor=MSO_ANCHOR.MIDDLE)
_txt(s, 14.5, 3.05, 12.0, 0.5, "공급 확대 4단", size=11, bold=True, color=NAVY2)
flow4 = ["① 중소형 Supermarket 중심 가공상품 공급",
         "② 대형 Supermarket 중심 PB·공산품 공급",
         "③ 신선식품 공급 및 마스터 프랜차이즈 집중 물색",
         "④ Master Franchise 계약 → 확장 집중"]
fy = 3.7
for i, ft in enumerate(flow4):
    flownode(s, 14.5, fy, 11.8, 1.9, ft, fill=(RED if i == 3 else NAVY2), size=12)
    if i < 3:
        arrow(s, 20.4, fy + 1.9, 20.4, fy + 2.35, color=NAVY2, w=2.5)
    fy += 2.35
pageno(s)

# ══════════ 4-4 향후 계획_제언 ══════════
s = newslide()
hdr(s, 4, "향후 계획_제언")
lead(s, "저자본 B2B 파일럿 우선 추진, Gate 충족 시 단계 확장")
box(s, X0, 3.6, 8.0, 10.0, "제언",
    "베트남 남부 인프라 활용한\nB2B 파일럿 先 추진\n\n→ 시장성 검증 후\n프랜차이즈 확장")
box(s, 9.56, 3.6, 8.0, 10.0, "기대효과",
    "① 신규 투자 최소\n   (기존 인프라 연장)\n\n② 진입 1년 내 현금흐름 창출\n\n"
    "③ 베트남 단일 의존 분산\n\n(정량 기대효과 : 〔조사〕)")
box(s, 17.92, 3.6, 8.4, 10.0, "향후 검토과제",
    "① D-8 시장조사로 Gate 수치 확정\n\n② 프랜차이즈 법제 제정 동향 모니터링\n\n"
    "③ 파트너 3사(Chip Mong·Royal Group·부영) 접촉 우선순위 수립", tfill=RED)
pageno(s)

out = os.path.join(os.path.dirname(__file__), "..", "output", "롯데마트_캄보디아_확장전략_v0.3.pptx")
prs.save(out)
print("saved", os.path.normpath(out), "| slides:", len(prs.slides._sldIdLst))
